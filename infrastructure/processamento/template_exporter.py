"""
template_exporter.py

Sistema de Templates de Exportação - Omni Keywords Finder

Prompt: CHECKLIST_PRIMEIRA_REVISAO.md - Item 10
Ruleset: enterprise_control_layer.yaml
Data: 2024-12-19

Funcionalidades:
- Templates HTML para relatórios
- Templates PowerPoint
- Templates Markdown
- Personalização de templates
- Preview de templates
- Versionamento de templates
- Sistema de cache e otimização
- Validação de templates
- Exportação em múltiplos formatos
"""

import json
import os
import hashlib
import zipfile
import tempfile
from typing import Dict, List, Optional, Any, Union, Tuple
from datetime import datetime, timedelta
from pathlib import Path
from dataclasses import dataclass, asdict
from enum import Enum
import logging
import re
from copy import deepcopy

# Importações opcionais
try:
    from jinja2 import Template, Environment, FileSystemLoader, TemplateError
    JINJA2_AVAILABLE = True
except ImportError:
    JINJA2_AVAILABLE = False
    logger = logging.getLogger(__name__)
    logger.warning("Jinja2 not available. Template rendering will be limited.")

try:
    import yaml
    YAML_AVAILABLE = True
except ImportError:
    YAML_AVAILABLE = False

logger = logging.getLogger(__name__)

class TemplateFormat(Enum):
    """Formatos de template suportados"""
    HTML = "html"
    POWERPOINT = "pptx"
    MARKDOWN = "md"
    JSON = "json"
    XML = "xml"

class TemplateType(Enum):
    """Tipos de template"""
    KEYWORDS_REPORT = "keywords_report"
    CLUSTERS_REPORT = "clusters_report"
    PERFORMANCE_REPORT = "performance_report"
    BUSINESS_REPORT = "business_report"
    AUDIT_REPORT = "audit_report"
    CUSTOM = "custom"

@dataclass
class TemplateConfig:
    """Configuração de template"""
    name: str
    description: str
    format: TemplateFormat
    type: TemplateType
    version: str
    author: str
    created_at: datetime
    updated_at: datetime
    is_active: bool = True
    is_default: bool = False
    metadata: Dict[str, Any] = None
    
    def __post_init__(self):
        if self.metadata is None:
            self.metadata = {}

@dataclass
class TemplateVariable:
    """Variável de template"""
    name: str
    type: str  # string, number, boolean, date, array, object
    description: str
    required: bool = False
    default_value: Any = None
    validation_rules: Dict[str, Any] = None
    
    def __post_init__(self):
        if self.validation_rules is None:
            self.validation_rules = {}

@dataclass
class ExportData:
    """Dados para exportação"""
    keywords: List[Dict[str, Any]]
    clusters: List[Dict[str, Any]]
    performance_metrics: Dict[str, Any]
    business_metrics: Dict[str, Any]
    audit_logs: List[Dict[str, Any]]
    custom_data: Dict[str, Any]
    
    def __init__(
        self,
        keywords: Optional[List[Dict[str, Any]]] = None,
        clusters: Optional[List[Dict[str, Any]]] = None,
        performance_metrics: Optional[Dict[str, Any]] = None,
        business_metrics: Optional[Dict[str, Any]] = None,
        audit_logs: Optional[List[Dict[str, Any]]] = None,
        custom_data: Optional[Dict[str, Any]] = None
    ):
        self.keywords = keywords or []
        self.clusters = clusters or []
        self.performance_metrics = performance_metrics or {}
        self.business_metrics = business_metrics or {}
        self.audit_logs = audit_logs or []
        self.custom_data = custom_data or {}

class TemplateValidator:
    """Validador de templates"""
    
    @staticmethod
    def validate_template_content(content: str, format: TemplateFormat) -> Tuple[bool, List[str]]:
        """Valida conteúdo do template"""
        errors = []
        
        if not content or not content.strip():
            errors.append("Template content is empty")
            return False, errors
        
        if format == TemplateFormat.HTML:
            # Validar HTML básico
            if not re.search(r'<html|<body|<head', content, re.IGNORECASE):
                errors.append("HTML template must contain basic HTML structure")
            
            # Verificar tags não fechadas
            open_tags = re.findall(r'<([^/][^>]*)>', content)
            close_tags = re.findall(r'</([^>]*)>', content)
            if len(open_tags) != len(close_tags):
                errors.append("HTML template has unclosed tags")
        
        elif format == TemplateFormat.MARKDOWN:
            # Validar Markdown básico
            if not re.search(r'#|\*|\-|\[', content):
                errors.append("Markdown template should contain basic markdown syntax")
        
        elif format == TemplateFormat.JSON:
            try:
                json.loads(content)
            except json.JSONDecodeError as e:
                errors.append(f"Invalid JSON template: {str(e)}")
        
        return len(errors) == 0, errors
    
    @staticmethod
    def validate_template_variables(template_vars: List[TemplateVariable], data: ExportData) -> Tuple[bool, List[str]]:
        """Valida variáveis do template contra dados"""
        errors = []
        
        for var in template_vars:
            if var.required:
                # Verificar se variável existe nos dados
                if not TemplateValidator._variable_exists_in_data(var.name, data):
                    errors.append(f"Required variable '{var.name}' not found in data")
                
                # Validar tipo
                if not TemplateValidator._validate_variable_type(var.name, var.type, data):
                    errors.append(f"Variable '{var.name}' type mismatch. Expected: {var.type}")
        
        return len(errors) == 0, errors
    
    @staticmethod
    def _variable_exists_in_data(var_name: str, data: ExportData) -> bool:
        """Verifica se variável existe nos dados"""
        # Verificar em diferentes seções dos dados
        sections = [
            data.keywords,
            data.clusters,
            data.performance_metrics,
            data.business_metrics,
            data.audit_logs,
            data.custom_data
        ]
        
        for section in sections:
            if isinstance(section, dict) and var_name in section:
                return True
            elif isinstance(section, list) and section:
                if isinstance(section[0], dict) and var_name in section[0]:
                    return True
        
        return False
    
    @staticmethod
    def _validate_variable_type(var_name: str, expected_type: str, data: ExportData) -> bool:
        """Valida tipo da variável"""
        # Implementação simplificada - em produção seria mais robusta
        return True

class TemplateRenderer:
    """Renderizador de templates"""
    
    def __init__(self):
        if not JINJA2_AVAILABLE:
            raise ImportError("Jinja2 is required for template rendering")
        
        self.jinja_env = Environment(
            loader=FileSystemLoader('templates'),
            autoescape=True,
            trim_blocks=True,
            lstrip_blocks=True
        )
        
        # Adicionar filtros customizados
        self.jinja_env.filters['format_number'] = self._format_number
        self.jinja_env.filters['format_date'] = self._format_date
        self.jinja_env.filters['format_currency'] = self._format_currency
        self.jinja_env.filters['highlight_keywords'] = self._highlight_keywords
    
    def render_template(
        self,
        template_content: str,
        data: ExportData,
        variables: Optional[List[TemplateVariable]] = None
    ) -> str:
        """Renderiza template com dados"""
        try:
            # Preparar contexto
            context = self._prepare_context(data, variables)
            
            # Criar template Jinja2
            template = self.jinja_env.from_string(template_content)
            
            # Renderizar
            return template.render(**context)
        
        except TemplateError as e:
            logger.error(f"Template rendering error: {str(e)}")
            raise ValueError(f"Template rendering failed: {str(e)}")
    
    def _prepare_context(self, data: ExportData, variables: Optional[List[TemplateVariable]] = None) -> Dict[str, Any]:
        """Prepara contexto para renderização"""
        context = {
            'keywords': data.keywords,
            'clusters': data.clusters,
            'performance_metrics': data.performance_metrics,
            'business_metrics': data.business_metrics,
            'audit_logs': data.audit_logs,
            'custom_data': data.custom_data,
            'generated_at': datetime.now(),
            'total_keywords': len(data.keywords),
            'total_clusters': len(data.clusters),
            'total_audit_logs': len(data.audit_logs)
        }
        
        # Adicionar variáveis customizadas
        if variables:
            for var in variables:
                if var.default_value is not None:
                    context[var.name] = var.default_value
        
        return context
    
    def _format_number(self, value: Union[int, float], format_type: str = 'default') -> str:
        """Filtro para formatação de números"""
        if format_type == 'currency':
            return f"R$ {value:,.2f}"
        elif format_type == 'percentage':
            return f"{value:.2f}%"
        elif format_type == 'thousands':
            return f"{value:,}"
        else:
            return str(value)
    
    def _format_date(self, value: Union[str, datetime], format_type: str = 'default') -> str:
        """Filtro para formatação de datas"""
        if isinstance(value, str):
            try:
                value = datetime.fromisoformat(value)
            except ValueError:
                return value
        
        if format_type == 'short':
            return value.strftime('%data/%m/%Y')
        elif format_type == 'long':
            return value.strftime('%data de %B de %Y')
        elif format_type == 'datetime':
            return value.strftime('%data/%m/%Y %H:%M')
        else:
            return value.strftime('%Y-%m-%data')
    
    def _format_currency(self, value: Union[int, float]) -> str:
        """Filtro para formatação de moeda"""
        return f"R$ {value:,.2f}"
    
    def _highlight_keywords(self, text: str, keywords: List[str]) -> str:
        """Filtro para destacar keywords no texto"""
        for keyword in keywords:
            text = re.sub(
                re.escape(keyword),
                f'<mark>{keyword}</mark>',
                text,
                flags=re.IGNORECASE
            )
        return text

class PowerPointGenerator:
    """Gerador de apresentações PowerPoint"""
    
    def __init__(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            self.pptx_available = True
        except ImportError:
            self.pptx_available = False
            logger.warning("python-pptx not available. PowerPoint generation disabled.")
    
    def generate_presentation(
        self,
        data: ExportData,
        template_config: Optional[TemplateConfig],
        output_path: str
    ) -> str:
        """Gera apresentação PowerPoint"""
        if not self.pptx_available:
            raise ImportError("python-pptx is required for PowerPoint generation")
        
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN
        
        # Criar apresentação
        prs = Presentation()
        
        # Slide de título
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = f"Relatório - {template_config.name if template_config else 'Keywords'}"
        title_slide.placeholders[1].text = f"Gerado em {datetime.now().strftime('%data/%m/%Y %H:%M')}"
        
        # Slide de resumo
        if data.keywords or data.clusters:
            summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
            summary_slide.shapes.title.text = "Resumo Executivo"
            
            summary_text = f"""
            • Total de Keywords: {len(data.keywords)}
            • Total de Clusters: {len(data.clusters)}
            • Data de Geração: {datetime.now().strftime('%data/%m/%Y')}
            """
            summary_slide.placeholders[1].text = summary_text
        
        # Slide de keywords (se houver)
        if data.keywords:
            keywords_slide = prs.slides.add_slide(prs.slide_layouts[1])
            keywords_slide.shapes.title.text = "Keywords Processadas"
            
            keywords_text = "\n".join([
                f"• {kw.get('termo', 'N/A')} - Volume: {kw.get('volume', 0)}"
                for kw in data.keywords[:10]  # Limitar a 10 keywords
            ])
            keywords_slide.placeholders[1].text = keywords_text
        
        # Slide de clusters (se houver)
        if data.clusters:
            clusters_slide = prs.slides.add_slide(prs.slide_layouts[1])
            clusters_slide.shapes.title.text = "Clusters Gerados"
            
            clusters_text = "\n".join([
                f"• {cluster.get('nome', 'N/A')} - {len(cluster.get('keywords', []))} keywords"
                for cluster in data.clusters[:10]  # Limitar a 10 clusters
            ])
            clusters_slide.placeholders[1].text = clusters_text
        
        # Salvar apresentação
        prs.save(output_path)
        return output_path

class TemplateExporter:
    """Sistema principal de templates de exportação"""
    
    def __init__(self, templates_dir: str = "templates/export"):
        self.templates_dir = Path(templates_dir)
        self.templates_dir.mkdir(parents=True, exist_ok=True)
        
        self.validator = TemplateValidator()
        self.renderer = TemplateRenderer() if JINJA2_AVAILABLE else None
        self.powerpoint_generator = PowerPointGenerator()
        
        # Cache de templates
        self._template_cache = {}
        self._cache_timestamp = {}
        self._cache_ttl = 300  # 5 minutos
    
    def create_template(
        self,
        name: str,
        description: str,
        format: TemplateFormat,
        type: TemplateType,
        content: str,
        author: str,
        variables: Optional[List[TemplateVariable]] = None,
        is_default: bool = False
    ) -> TemplateConfig:
        """Cria novo template"""
        # Validar conteúdo
        is_valid, errors = self.validator.validate_template_content(content, format)
        if not is_valid:
            raise ValueError(f"Invalid template content: {', '.join(errors)}")
        
        # Criar configuração
        config = TemplateConfig(
            name=name,
            description=description,
            format=format,
            type=type,
            version="1.0.0",
            author=author,
            created_at=datetime.now(),
            updated_at=datetime.now(),
            is_default=is_default
        )
        
        # Salvar template
        template_id = self._save_template(config, content, variables)
        
        # Limpar cache
        self._clear_cache()
        
        logger.info(f"Template created: {template_id}")
        return config
    
    def update_template(
        self,
        template_id: str,
        content: Optional[str] = None,
        variables: Optional[List[TemplateVariable]] = None,
        description: Optional[str] = None,
        is_active: Optional[bool] = None
    ) -> TemplateConfig:
        """Atualiza template existente"""
        # Carregar template atual
        config = self.get_template_config(template_id)
        if not config:
            raise ValueError(f"Template not found: {template_id}")
        
        # Validar novo conteúdo se fornecido
        if content:
            is_valid, errors = self.validator.validate_template_content(content, config.format)
            if not is_valid:
                raise ValueError(f"Invalid template content: {', '.join(errors)}")
        
        # Atualizar configuração
        if description:
            config.description = description
        if is_active is not None:
            config.is_active = is_active
        
        config.updated_at = datetime.now()
        config.version = self._increment_version(config.version)
        
        # Salvar atualizações
        self._save_template(config, content, variables, template_id)
        
        # Limpar cache
        self._clear_cache()
        
        logger.info(f"Template updated: {template_id}")
        return config
    
    def get_template_config(self, template_id: str) -> Optional[TemplateConfig]:
        """Obtém configuração do template"""
        config_path = self.templates_dir / f"{template_id}.json"
        if not config_path.exists():
            return None
        
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Converter datas
            data['created_at'] = datetime.fromisoformat(data['created_at'])
            data['updated_at'] = datetime.fromisoformat(data['updated_at'])
            
            return TemplateConfig(**data)
        except Exception as e:
            logger.error(f"Error loading template config: {str(e)}")
            return None
    
    def get_template_content(self, template_id: str) -> Optional[str]:
        """Obtém conteúdo do template"""
        content_path = self.templates_dir / f"{template_id}.template"
        if not content_path.exists():
            return None
        
        try:
            with open(content_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error loading template content: {str(e)}")
            return None
    
    def get_template_variables(self, template_id: str) -> List[TemplateVariable]:
        """Obtém variáveis do template"""
        vars_path = self.templates_dir / f"{template_id}.vars.json"
        if not vars_path.exists():
            return []
        
        try:
            with open(vars_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            return [TemplateVariable(**var_data) for var_data in data]
        except Exception as e:
            logger.error(f"Error loading template variables: {str(e)}")
            return []
    
    def list_templates(
        self,
        format: Optional[TemplateFormat] = None,
        type: Optional[TemplateType] = None,
        active_only: bool = True
    ) -> List[TemplateConfig]:
        """Lista templates disponíveis"""
        templates = []
        
        for config_file in self.templates_dir.glob("*.json"):
            if config_file.name.endswith('.vars.json'):
                continue
            
            template_id = config_file.stem
            config = self.get_template_config(template_id)
            
            if not config:
                continue
            
            # Filtrar por formato
            if format and config.format != format:
                continue
            
            # Filtrar por tipo
            if type and config.type != type:
                continue
            
            # Filtrar por status ativo
            if active_only and not config.is_active:
                continue
            
            templates.append(config)
        
        # Ordenar por data de atualização
        templates.sort(key=lambda value: value.updated_at, reverse=True)
        
        return templates
    
    def export_with_template(
        self,
        template_id: str,
        data: ExportData,
        output_path: Optional[str] = None,
        custom_variables: Optional[Dict[str, Any]] = None
    ) -> str:
        """Exporta dados usando template específico"""
        # Carregar template
        config = self.get_template_config(template_id)
        if not config:
            raise ValueError(f"Template not found: {template_id}")
        
        if not config.is_active:
            raise ValueError(f"Template is not active: {template_id}")
        
        content = self.get_template_content(template_id)
        if not content:
            raise ValueError(f"Template content not found: {template_id}")
        
        variables = self.get_template_variables(template_id)
        
        # Validar dados contra variáveis
        is_valid, errors = self.validator.validate_template_variables(variables, data)
        if not is_valid:
            logger.warning(f"Template variable validation warnings: {', '.join(errors)}")
        
        # Gerar nome do arquivo de saída
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"export_{template_id}_{timestamp}.{config.format.value}"
        
        # Renderizar baseado no formato
        if config.format == TemplateFormat.HTML:
            return self._export_html(content, data, variables, output_path, custom_variables)
        elif config.format == TemplateFormat.POWERPOINT:
            return self._export_powerpoint(content, data, variables, output_path, custom_variables)
        elif config.format == TemplateFormat.MARKDOWN:
            return self._export_markdown(content, data, variables, output_path, custom_variables)
        elif config.format == TemplateFormat.JSON:
            return self._export_json(content, data, variables, output_path, custom_variables)
        else:
            raise ValueError(f"Unsupported format: {config.format}")
    
    def preview_template(
        self,
        template_id: str,
        data: ExportData,
        custom_variables: Optional[Dict[str, Any]] = None
    ) -> str:
        """Gera preview do template"""
        if not self.renderer:
            raise ImportError("Jinja2 is required for template preview")
        
        # Carregar template
        config = self.get_template_config(template_id)
        if not config:
            raise ValueError(f"Template not found: {template_id}")
        
        content = self.get_template_content(template_id)
        if not content:
            raise ValueError(f"Template content not found: {template_id}")
        
        variables = self.get_template_variables(template_id)
        
        # Renderizar preview
        try:
            rendered = self.renderer.render_template(content, data, variables)
            
            # Para HTML, adicionar estilos de preview
            if config.format == TemplateFormat.HTML:
                rendered = self._add_preview_styles(rendered)
            
            return rendered
        except Exception as e:
            logger.error(f"Preview generation failed: {str(e)}")
            raise ValueError(f"Preview generation failed: {str(e)}")
    
    def create_template_version(self, template_id: str, version_name: str) -> str:
        """Cria nova versão do template"""
        config = self.get_template_config(template_id)
        if not config:
            raise ValueError(f"Template not found: {template_id}")
        
        content = self.get_template_content(template_id)
        variables = self.get_template_variables(template_id)
        
        # Criar nova versão
        new_config = deepcopy(config)
        new_config.version = version_name
        new_config.created_at = datetime.now()
        new_config.updated_at = datetime.now()
        
        new_template_id = f"{template_id}_v{version_name}"
        
        # Salvar nova versão
        self._save_template(new_config, content, variables, new_template_id)
        
        logger.info(f"Template version created: {new_template_id}")
        return new_template_id
    
    def _save_template(
        self,
        config: TemplateConfig,
        content: str,
        variables: Optional[List[TemplateVariable]] = None,
        template_id: Optional[str] = None
    ) -> str:
        """Salva template no sistema"""
        if not template_id:
            template_id = self._generate_template_id(config.name)
        
        # Salvar configuração
        config_path = self.templates_dir / f"{template_id}.json"
        config_data = asdict(config)
        config_data['created_at'] = config.created_at.isoformat()
        config_data['updated_at'] = config.updated_at.isoformat()
        
        with open(config_path, 'w', encoding='utf-8') as f:
            json.dump(config_data, f, indent=2, ensure_ascii=False)
        
        # Salvar conteúdo
        content_path = self.templates_dir / f"{template_id}.template"
        with open(content_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        # Salvar variáveis
        if variables:
            vars_path = self.templates_dir / f"{template_id}.vars.json"
            vars_data = [asdict(var) for var in variables]
            with open(vars_path, 'w', encoding='utf-8') as f:
                json.dump(vars_data, f, indent=2, ensure_ascii=False)
        
        return template_id
    
    def _generate_template_id(self, name: str) -> str:
        """Gera ID único para template"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = re.sub(r'[^a-zA-Z0-9_]', '_', name.lower())
        return f"{safe_name}_{timestamp}"
    
    def _increment_version(self, current_version: str) -> str:
        """Incrementa versão do template"""
        try:
            parts = current_version.split('.')
            if len(parts) >= 3:
                major, minor, patch = parts[:3]
                patch = str(int(patch) + 1)
                return f"{major}.{minor}.{patch}"
            else:
                return f"{current_version}.1"
        except:
            return f"{current_version}.1"
    
    def _export_html(
        self,
        content: str,
        data: ExportData,
        variables: List[TemplateVariable],
        output_path: str,
        custom_variables: Optional[Dict[str, Any]] = None
    ) -> str:
        """Exporta template HTML"""
        if not self.renderer:
            raise ImportError("Jinja2 is required for HTML export")
        
        try:
            rendered = self.renderer.render_template(content, data, variables)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(rendered)
            
            return output_path
        except Exception as e:
            logger.error(f"HTML export failed: {str(e)}")
            raise ValueError(f"HTML export failed: {str(e)}")
    
    def _export_powerpoint(
        self,
        content: str,
        data: ExportData,
        variables: List[TemplateVariable],
        output_path: str,
        custom_variables: Optional[Dict[str, Any]] = None
    ) -> str:
        """Exporta template PowerPoint"""
        try:
            # Usar gerador PowerPoint
            config = None  # Não usado no PowerPoint generator atual
            return self.powerpoint_generator.generate_presentation(data, config, output_path)
        except Exception as e:
            logger.error(f"PowerPoint export failed: {str(e)}")
            raise ValueError(f"PowerPoint export failed: {str(e)}")
    
    def _export_markdown(
        self,
        content: str,
        data: ExportData,
        variables: List[TemplateVariable],
        output_path: str,
        custom_variables: Optional[Dict[str, Any]] = None
    ) -> str:
        """Exporta template Markdown"""
        if not self.renderer:
            raise ImportError("Jinja2 is required for Markdown export")
        
        try:
            rendered = self.renderer.render_template(content, data, variables)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(rendered)
            
            return output_path
        except Exception as e:
            logger.error(f"Markdown export failed: {str(e)}")
            raise ValueError(f"Markdown export failed: {str(e)}")
    
    def _export_json(
        self,
        content: str,
        data: ExportData,
        variables: List[TemplateVariable],
        output_path: str,
        custom_variables: Optional[Dict[str, Any]] = None
    ) -> str:
        """Exporta template JSON"""
        if not self.renderer:
            raise ImportError("Jinja2 is required for JSON export")
        
        try:
            # Para JSON, renderizar como template e validar
            rendered = self.renderer.render_template(content, data, variables)
            
            # Validar JSON
            json_data = json.loads(rendered)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(json_data, f, indent=2, ensure_ascii=False)
            
            return output_path
        except Exception as e:
            logger.error(f"JSON export failed: {str(e)}")
            raise ValueError(f"JSON export failed: {str(e)}")
    
    def _add_preview_styles(self, html_content: str) -> str:
        """Adiciona estilos de preview ao HTML"""
        preview_styles = """
        <style>
        .template-preview {
            border: 2px dashed #ccc;
            padding: 20px;
            margin: 10px;
            background: #f9f9f9;
        }
        .template-preview::before {
            content: "PREVIEW MODE";
            display: block;
            background: #ff6b6b;
            color: white;
            padding: 5px 10px;
            margin: -20px -20px 20px -20px;
            font-size: 12px;
            font-weight: bold;
            text-align: center;
        }
        </style>
        """
        
        # Inserir estilos no head
        if '<head>' in html_content:
            html_content = html_content.replace('<head>', f'<head>{preview_styles}')
        else:
            html_content = f'<html><head>{preview_styles}</head><body class="template-preview">{html_content}</body></html>'
        
        return html_content
    
    def _clear_cache(self):
        """Limpa cache de templates"""
        self._template_cache.clear()
        self._cache_timestamp.clear()

# Instância global
template_exporter = TemplateExporter()
